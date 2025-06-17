import os
import re
from typing import List, Dict, Any, Optional
from pathlib import Path
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import magic

class DocumentProcessor:
    """
    Classe pour traiter différents types de documents (PDF, TXT, DOCX, PPTX).
    """
    
    def __init__(self, upload_dir: str = "data/uploads"):
        """
        Initialise le processeur de documents.
        
        Args:
            upload_dir (str): Répertoire où les documents sont stockés
        """
        self.upload_dir = Path(upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self.mime = magic.Magic(mime=True)
    
    def get_file_type(self, file_path: str) -> str:
        """Détermine le type de fichier."""
        return self.mime.from_file(file_path)
    
    def process_document(self, file_path: str) -> Dict[str, Any]:
        """
        Traite un document et extrait son contenu texte.
        
        Args:
            file_path (str): Chemin vers le fichier à traiter
            
        Returns:
            dict: Dictionnaire contenant le contenu texte et les métadonnées
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"Le fichier {file_path} n'existe pas.")
        
        file_type = self.get_file_type(str(file_path))
        
        if 'pdf' in file_type:
            return self._process_pdf(file_path)
        elif 'text/plain' in file_type or file_path.suffix.lower() == '.txt':
            return self._process_txt(file_path)
        elif 'word' in file_type or file_path.suffix.lower() in ['.docx', '.doc']:
            return self._process_docx(file_path)
        elif 'presentation' in file_type or file_path.suffix.lower() in ['.pptx', '.ppt']:
            return self._process_pptx(file_path)
        else:
            raise ValueError(f"Type de fichier non supporté: {file_type}")
    
    def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Traite un fichier PDF et extrait le texte."""
        text = ""
        metadata = {
            "title": file_path.stem,
            "type": "pdf",
            "pages": 0,
            "sections": []
        }
        
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            metadata["pages"] = len(pdf_reader.pages)
            
            for i, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text += f"\n\n--- Page {i} ---\n{page_text}"
                    metadata["sections"].append({"page": i, "start": len(text) - len(page_text)})
        
        return {
            "content": text.strip(),
            "metadata": metadata
        }
    
    def _process_txt(self, file_path: Path) -> Dict[str, Any]:
        """Traite un fichier texte."""
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        return {
            "content": content,
            "metadata": {
                "title": file_path.stem,
                "type": "text",
                "sections": [{"section": 1, "start": 0}]
            }
        }
    
    def _process_docx(self, file_path: Path) -> Dict[str, Any]:
        """Traite un document Word."""
        doc = docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        content = "\n\n".join(paragraphs)
        
        return {
            "content": content,
            "metadata": {
                "title": file_path.stem,
                "type": "docx",
                "sections": [{"heading": p.text.strip(), "start": 0} for i, p in enumerate(doc.paragraphs) 
                             if p.style.name.startswith('Heading')]
            }
        }
    
    def _process_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Traite une présentation PowerPoint."""
        prs = Presentation(file_path)
        content = []
        metadata = {
            "title": file_path.stem,
            "type": "pptx",
            "slides": []
        }
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            slide_content = "\n".join(filter(None, slide_text))
            if slide_content:
                content.append(f"--- Slide {i} ---\n{slide_content}")
                metadata["slides"].append({"number": i, "start": len("\n".join(content)) - len(slide_content)})
        
        return {
            "content": "\n\n".join(content),
            "metadata": metadata
        }
    
    def clean_text(self, text: str) -> str:
        """Nettoie le texte en supprimant les caractères spéciaux et les espaces superflus."""
        # Supprimer les sauts de ligne multiples
        text = re.sub(r'\n+', '\n', text)
        # Supprimer les espaces multiples
        text = re.sub(r'\s+', ' ', text)
        # Supprimer les caractères non imprimables
        text = re.sub(r'[\x00-\x1f\x7f-\x9f]', '', text)
        return text.strip()
    
    def split_into_chunks(self, text: str, chunk_size: int = 1000, overlap: int = 100) -> List[Dict[str, Any]]:
        """
        Découpe le texte en morceaux avec un chevauchement.
        
        Args:
            text (str): Texte à découper
            chunk_size (int): Taille maximale d'un morceau
            overlap (int): Nombre de caractères de chevauchement entre les morceaux
            
        Returns:
            List[Dict[str, Any]]: Liste de dictionnaires contenant les morceaux de texte et leurs métadonnées
        """
        if not text:
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = min(start + chunk_size, text_length)
            
            # Ajuster la fin pour ne pas couper les mots
            if end < text_length:
                # Trouver le dernier espace avant la fin
                last_space = text.rfind(' ', start, end + 1)
                if last_space > start and (last_space - start) > (chunk_size // 2):
                    end = last_space
            
            chunk = text[start:end].strip()
            if chunk:  # Ne pas ajouter de morceaux vides
                chunks.append({
                    "text": chunk,
                    "start": start,
                    "end": end,
                    "length": len(chunk)
                })
            
            # Déplacer le point de départ pour le prochain morceau
            if end == text_length:
                break
                
            start = end - overlap
            if start <= 0:
                start = end
        
        return chunks
